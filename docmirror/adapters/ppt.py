"""
PPT Adapter — .pptx → BaseResult
==================================

Extracts slide content from PowerPoint presentations using python-pptx.

Processing logic:
    1. Opens the .pptx file via ``python-pptx.Presentation``.
    2. Each slide becomes a separate ``PageLayout`` (page index = slide order).
    3. For each slide:
       - If the slide has a title shape with text, a ``title`` Block is
         created (heading level 2) and used as the slide's Markdown heading.
       - All other shapes are iterated:
         * Text-bearing shapes → ``text`` Block
         * Table shapes → ``table`` Block (each row as a list of cell text)
       - The slide's title shape is skipped during the iteration to avoid
         duplicate content.
    4. The full_text joins all slide content separated by double newlines,
       with each slide prefixed by "### Slide N: {title}".

Metadata includes:
    - source_format: "pptx"
    - slide_count: total number of slides
"""

from __future__ import annotations

import logging
from pathlib import Path

from docmirror.framework.base import BaseParser, ParserOutput, ParserStatus
from docmirror.models.domain import BaseResult, Block, PageLayout

logger = logging.getLogger(__name__)


class PPTAdapter(BaseParser):
    """PowerPoint (.pptx) format adapter — each slide becomes a separate page."""

    async def to_base_result(self, file_path: Path) -> BaseResult:
        """
        Parse a .pptx file into a BaseResult.

        Each slide maps to a PageLayout. Slide titles become title Blocks,
        text shapes become text Blocks, and table shapes become table Blocks.
        """
        from pptx import Presentation
        prs = Presentation(str(file_path))

        pages = []
        text_parts = []

        for i, slide in enumerate(prs.slides):
            blocks = []
            slide_texts = []

            # Extract slide title (if present)
            if slide.shapes.title and slide.shapes.title.text:
                blocks.append(Block(
                    block_type="title",
                    raw_content=slide.shapes.title.text,
                    page=i,
                    heading_level=2,
                ))
                slide_texts.append(f"### Slide {i+1}: {slide.shapes.title.text}")
            else:
                slide_texts.append(f"### Slide {i+1}")

            # Extract content from all other shapes (skip the title shape)
            for shape in slide.shapes:
                if shape == slide.shapes.title:
                    continue

                # Text content from text frames
                if hasattr(shape, "text") and shape.text:
                    blocks.append(Block(block_type="text", raw_content=shape.text, page=i))
                    slide_texts.append(shape.text)

                # Table content — each row becomes a list of cell text values
                if shape.has_table:
                    rows = [[cell.text.strip() for cell in row.cells] for row in shape.table.rows]
                    if rows:
                        blocks.append(Block(block_type="table", raw_content=rows, page=i))

            pages.append(PageLayout(page_number=i, blocks=tuple(blocks)))
            text_parts.append("\n".join(slide_texts))

        return BaseResult(
            pages=tuple(pages),
            full_text="\n\n".join(text_parts),
            metadata={"source_format": "pptx", "slide_count": len(prs.slides)},
        )
