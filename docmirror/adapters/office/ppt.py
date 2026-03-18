# Copyright (c) 2026 ValueMap Global and contributors. All rights reserved.
# Author: Adam Lin <adamlin@valuemapglobal.com>
#
# This source code is licensed under the Apache 2.0 license found in the
# LICENSE file in the root directory of this source tree.

"""
PPT Adapter — .pptx → ParseResult
==================================

Extracts slide content from PowerPoint presentations using python-pptx.

Processing logic:
    1. Opens the .pptx file via ``python-pptx.Presentation``.
    2. Each slide becomes a separate ``PageContent`` (page index = slide order).
    3. For each slide:
       - Slide title → ``TextBlock(level=H2)``
       - Text shapes → ``TextBlock(level=BODY)``
       - Table shapes → ``TableBlock`` with typed ``CellValue``

Metadata includes:
    - parser_name: "PPTAdapter"
    - page_count: total number of slides
"""

from __future__ import annotations

import logging
from pathlib import Path
from typing import List

from docmirror.framework.base import BaseParser

logger = logging.getLogger(__name__)


class PPTAdapter(BaseParser):
    """PowerPoint (.pptx/.ppt) format adapter — native parsing prioritizing `python-pptx`.

    Processing strategy:
        1. **Modern formats (.pptx):** Direct deep extraction via ``python-pptx``
           preserving slides, shapes, text, and tables natively.
        2. **Legacy formats (.ppt):** Automatically transcoded to PDF via LibreOffice,
           then processed through the standard PDF pipeline.
    """

    async def to_parse_result(self, file_path: Path, **kwargs) -> ParseResult:
        """
        Parse a .pptx file into a ParseResult.

        Each slide maps to a PageContent. Slide titles become H2 TextBlocks,
        text shapes become BODY TextBlocks, and table shapes become TableBlocks.
        """
        from pptx import Presentation

        from docmirror.models.entities.parse_result import (
            CellValue,
            DataType,
            PageContent,
            ParseResult,
            ParserInfo,
            TableBlock,
            TableRow,
            TextBlock,
            TextLevel,
        )

        logger.info(f"[PPTAdapter] Starting native extraction for presentation: {file_path}")
        prs = Presentation(str(file_path))

        pages: list[PageContent] = []

        for i, slide in enumerate(prs.slides):
            texts: list[TextBlock] = []
            tables: list[TableBlock] = []

            # Extract slide title (if present)
            if slide.shapes.title and slide.shapes.title.text:
                texts.append(
                    TextBlock(
                        content=slide.shapes.title.text,
                        level=TextLevel.H2,
                    )
                )

            # Extract content from all other shapes (skip the title shape)
            for shape in slide.shapes:
                if shape == slide.shapes.title:
                    continue

                # Text content from text frames
                if hasattr(shape, "text") and shape.text:
                    texts.append(TextBlock(content=shape.text, level=TextLevel.BODY))

                # Table content — typed CellValue
                if shape.has_table:
                    first_row = True
                    headers: list[str] = []
                    data_rows: list[TableRow] = []

                    for row in shape.table.rows:
                        cells = [CellValue(text=cell.text.strip(), data_type=DataType.TEXT) for cell in row.cells]
                        if first_row:
                            headers = [c.text for c in cells]
                            first_row = False
                        else:
                            if any(c.text for c in cells):
                                data_rows.append(TableRow(cells=cells, source_page=i))

                    if headers or data_rows:
                        tables.append(
                            TableBlock(
                                table_id=f"slide{i}_table{len(tables)}",
                                headers=headers,
                                rows=data_rows,
                                page=i,
                            )
                        )

            pages.append(PageContent(page_number=i, texts=texts, tables=tables))

        return ParseResult(
            pages=pages,
            parser_info=ParserInfo(
                parser_name="PPTAdapter",
                page_count=len(prs.slides),
                overall_confidence=1.0,
            ),
        )
